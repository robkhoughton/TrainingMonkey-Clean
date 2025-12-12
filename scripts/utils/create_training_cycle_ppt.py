from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Add blank slide
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# Add purple gradient background
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(120, 110, 200)  # Purple color

# Add title "The Training Cycle:"
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "The Training Cycle:"
title_para = title_frame.paragraphs[0]
title_para.alignment = PP_ALIGN.CENTER
title_para.font.size = Pt(44)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(255, 255, 255)

# Add subtitle
subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.4))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "How Your Training Monkey works to optimize your training and prevent injuries"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.alignment = PP_ALIGN.CENTER
subtitle_para.font.size = Pt(16)
subtitle_para.font.color.rgb = RGBColor(255, 255, 255)

# Center circle coordinates
center_x = Inches(5)
center_y = Inches(4)
circle_radius = Inches(1.2)

# Add center circle
center_circle = slide.shapes.add_shape(
    MSO_SHAPE.OVAL,
    center_x - circle_radius,
    center_y - circle_radius,
    circle_radius * 2,
    circle_radius * 2
)
center_circle.fill.solid()
center_circle.fill.fore_color.rgb = RGBColor(100, 90, 180)
center_circle.line.color.rgb = RGBColor(200, 200, 255)
center_circle.line.width = Pt(2)

# Function to add rounded rectangle box
def add_cycle_box(slide, x, y, width, height, title, subtitle, color):
    # Add rounded rectangle
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x, y, width, height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = RGBColor(200, 200, 255)
    box.line.width = Pt(2)

    # Add text
    text_frame = box.text_frame
    text_frame.clear()
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Title
    p = text_frame.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Subtitle
    p2 = text_frame.add_paragraph()
    p2.text = subtitle
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(240, 240, 255)

# Add the four cycle boxes
box_color = RGBColor(110, 100, 190)

# Train (top)
add_cycle_box(
    slide,
    Inches(3.8), Inches(2.1),
    Inches(2.4), Inches(0.9),
    "Train (record)",
    "Complete your workout\non Strava",
    box_color
)

# Monitor (right)
add_cycle_box(
    slide,
    Inches(6.5), Inches(3.5),
    Inches(2.4), Inches(0.9),
    "Monitor (past)",
    "Check Dashboard for\nyour current training\nstatus",
    box_color
)

# Journal (bottom)
add_cycle_box(
    slide,
    Inches(3.8), Inches(5.3),
    Inches(2.4), Inches(0.9),
    "Journal\n(present)",
    "Use the Journal to\nreflect on today's\nworkout",
    box_color
)

# Plan (left)
add_cycle_box(
    slide,
    Inches(1.2), Inches(3.5),
    Inches(2.4), Inches(0.9),
    "Plan (future)",
    "Get YTM Coach\nrecommendations for\nyour upcoming training\nstrategy",
    box_color
)

# Add curved arrows (using connectors)
# Top to Right
arrow1 = slide.shapes.add_connector(
    1,  # Straight connector
    Inches(6.2), Inches(2.6),
    Inches(7.2), Inches(3.7)
)
arrow1.line.color.rgb = RGBColor(200, 200, 255)
arrow1.line.width = Pt(3)

# Right to Bottom
arrow2 = slide.shapes.add_connector(
    1,
    Inches(7.2), Inches(4.6),
    Inches(6.2), Inches(5.5)
)
arrow2.line.color.rgb = RGBColor(200, 200, 255)
arrow2.line.width = Pt(3)

# Bottom to Left
arrow3 = slide.shapes.add_connector(
    1,
    Inches(3.8), Inches(5.5),
    Inches(2.8), Inches(4.6)
)
arrow3.line.color.rgb = RGBColor(200, 200, 255)
arrow3.line.width = Pt(3)

# Left to Top
arrow4 = slide.shapes.add_connector(
    1,
    Inches(2.8), Inches(3.7),
    Inches(3.8), Inches(2.6)
)
arrow4.line.color.rgb = RGBColor(200, 200, 255)
arrow4.line.width = Pt(3)

# Add "Quick Start Guide" section at bottom
qsg_title_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
qsg_title_frame = qsg_title_box.text_frame
qsg_title_frame.text = "Quick Start Guide"
qsg_para = qsg_title_frame.paragraphs[0]
qsg_para.alignment = PP_ALIGN.CENTER
qsg_para.font.size = Pt(32)
qsg_para.font.bold = True
qsg_para.font.color.rgb = RGBColor(40, 80, 150)

# Add numbered circles for steps
def add_step_circle(slide, x, y, number):
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        x, y,
        Inches(0.5), Inches(0.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(100, 90, 180)
    circle.line.color.rgb = RGBColor(100, 90, 180)

    text_frame = circle.text_frame
    text_frame.clear()
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = text_frame.paragraphs[0]
    p.text = str(number)
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

# Add step circles (partially visible in screenshot)
add_step_circle(slide, Inches(2.5), Inches(7), 1)
add_step_circle(slide, Inches(4.75), Inches(7), 2)
add_step_circle(slide, Inches(7), Inches(7), 3)

# Save presentation
prs.save('Training_Cycle_Presentation.pptx')
print("PowerPoint presentation created successfully: Training_Cycle_Presentation.pptx")
